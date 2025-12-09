#!/usr/bin/env python3
"""
Stage 1: PPT Data Extraction
Extracts structured content from PowerPoint presentations.
"""

import json
import sys
from pathlib import Path
from pptx import Presentation


def extract_slide_content(slide, slide_number):
    """
    Extract title and body content from a single slide.

    Args:
        slide: pptx.slide.Slide object
        slide_number: int, slide index (1-based)

    Returns:
        dict with slide data
    """
    slide_data = {
        "id": slide_number,
        "title": "",
        "content": "",
        "original_img": f"slide{slide_number}.png"
    }

    # Extract title
    if slide.shapes.title:
        slide_data["title"] = slide.shapes.title.text.strip()

    # Extract all text content from text boxes (ignoring positioning)
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            # Skip if this is the title (already captured)
            if shape == slide.shapes.title:
                continue
            text_parts.append(shape.text.strip())

    # Merge all text into one body
    slide_data["content"] = "\n".join(text_parts)

    return slide_data


def extract_presentation(ppt_path):
    """
    Extract all slides from a PowerPoint presentation.

    Args:
        ppt_path: str or Path, path to .pptx file

    Returns:
        list of slide dictionaries
    """
    ppt_path = Path(ppt_path)

    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT file not found: {ppt_path}")

    print(f"📖 Reading presentation: {ppt_path.name}")
    prs = Presentation(str(ppt_path))

    slides_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"  ├─ Extracting Slide {idx}...")
        slide_data = extract_slide_content(slide, idx)
        slides_data.append(slide_data)

    print(f"✅ Extracted {len(slides_data)} slides")
    return slides_data


def save_json(data, output_path):
    """Save data as formatted JSON."""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

    print(f"💾 Saved to: {output_path}")


def main():
    """Main entry point."""
    if len(sys.argv) < 2:
        print("Usage: python3 extract_ppt.py <path_to_ppt>")
        print("Example: python3 extract_ppt.py slide/survey-phase2-eng-PPT\\ \\(3\\).pptx")
        sys.exit(1)

    ppt_path = sys.argv[1]

    # Extract content
    slides_data = extract_presentation(ppt_path)

    # Save to JSON
    output_path = Path("output") / "extracted_slides.json"
    save_json(slides_data, output_path)

    # Print summary
    print("\n📊 Summary:")
    for slide in slides_data[:3]:  # Show first 3 slides
        print(f"\nSlide {slide['id']}: {slide['title']}")
        preview = slide['content'][:100] + "..." if len(slide['content']) > 100 else slide['content']
        print(f"  Content preview: {preview}")

    if len(slides_data) > 3:
        print(f"\n... and {len(slides_data) - 3} more slides")


if __name__ == "__main__":
    main()
