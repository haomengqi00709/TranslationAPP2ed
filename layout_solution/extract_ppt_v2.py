#!/usr/bin/env python3
"""
Stage 1 V2: PPT Data Extraction with Table Support
Extracts structured content from PowerPoint presentations, preserving tables.
"""

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from PIL import Image


def extract_table_data(table):
    """
    Extract table as structured data.

    Args:
        table: pptx table object

    Returns:
        dict with headers and rows
    """
    rows_data = []
    for row in table.rows:
        row_cells = []
        for cell in row.cells:
            row_cells.append(cell.text.strip())
        rows_data.append(row_cells)

    # Assume first row is headers
    if rows_data:
        return {
            "headers": rows_data[0],
            "rows": rows_data[1:]
        }
    return {"headers": [], "rows": []}


def extract_chart_info(shape, slide_number, chart_index):
    """
    Extract metadata about a chart shape.

    Args:
        shape: pptx shape object with chart
        slide_number: int, slide index
        chart_index: int, chart number on this slide

    Returns:
        dict with chart metadata
    """
    chart_info = {
        "chart_id": f"slide_{slide_number}_chart_{chart_index}",
        "chart_title": "",
        "chart_type": str(shape.chart.chart_type) if hasattr(shape, 'chart') else "unknown"
    }

    # Try to get chart title
    if hasattr(shape.chart, 'chart_title') and shape.chart.has_title:
        try:
            chart_info["chart_title"] = shape.chart.chart_title.text_frame.text
        except:
            pass

    return chart_info


def classify_slide_type(slide):
    """
    Determine what type of slide this is.

    Returns: "chart", "table", "text", "section_header"
    """
    charts = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.CHART]
    tables = [shape for shape in slide.shapes if shape.has_table]
    text_shapes = [shape for shape in slide.shapes
                   if hasattr(shape, 'text') and shape.text.strip()
                   and not shape.has_table and shape != slide.shapes.title]

    if charts:
        return "chart"
    elif tables:
        return "table"
    elif len(text_shapes) <= 2:  # Just title + subtitle or minimal text
        return "section_header"
    else:
        return "text"


def extract_slide_content(slide, slide_number):
    """
    Extract title, content, tables, and charts from a single slide.

    Args:
        slide: pptx.slide.Slide object
        slide_number: int, slide index (1-based)

    Returns:
        dict with slide data
    """
    slide_type = classify_slide_type(slide)

    slide_data = {
        "id": slide_number,
        "type": slide_type,
        "title": "",
        "content": "",
        "tables": [],
        "charts": [],
        "original_img": f"slide{slide_number}.png"
    }

    # Extract title
    if slide.shapes.title:
        slide_data["title"] = slide.shapes.title.text.strip()

    # Extract tables if present
    table_shapes = [shape for shape in slide.shapes if shape.has_table]
    if table_shapes:
        for table_shape in table_shapes:
            table_data = extract_table_data(table_shape.table)
            slide_data["tables"].append(table_data)

    # Extract charts if present
    chart_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.CHART]
    if chart_shapes:
        for idx, chart_shape in enumerate(chart_shapes, start=1):
            chart_info = extract_chart_info(chart_shape, slide_number, idx)
            slide_data["charts"].append(chart_info)

    # Extract text content (excluding tables, charts, and title)
    text_parts = []
    for shape in slide.shapes:
        # Skip title
        if shape == slide.shapes.title:
            continue
        # Skip tables
        if shape.has_table:
            continue
        # Skip charts
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            continue
        # Get text
        if hasattr(shape, "text") and shape.text.strip():
            text_parts.append(shape.text.strip())

    # Merge all text into one body
    slide_data["content"] = "\n".join(text_parts)

    return slide_data


def extract_presentation(ppt_path):
    """
    Extract all slides from a PowerPoint presentation.

    Args:
        ppt_path: str or Path, path to .pptx file

    Returns:
        list of slide dictionaries
    """
    ppt_path = Path(ppt_path)

    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT file not found: {ppt_path}")

    print(f"📖 Reading presentation: {ppt_path.name}")
    prs = Presentation(str(ppt_path))

    slides_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_data = extract_slide_content(slide, idx)

        # Print type indicator
        type_icon = {
            "chart": "📈",
            "table": "📊",
            "text": "📝",
            "section_header": "📌"
        }
        icon = type_icon.get(slide_data["type"], "❓")

        print(f"  ├─ Slide {idx} {icon} [{slide_data['type']}]", end="")
        if slide_data["charts"]:
            print(f" ({len(slide_data['charts'])} chart(s))", end="")
        if slide_data["tables"]:
            print(f" ({len(slide_data['tables'])} table(s))", end="")
        print()

        slides_data.append(slide_data)

    print(f"✅ Extracted {len(slides_data)} slides")

    # Summary
    type_counts = {}
    total_charts = 0
    total_tables = 0
    for slide in slides_data:
        slide_type = slide['type']
        type_counts[slide_type] = type_counts.get(slide_type, 0) + 1
        total_charts += len(slide.get('charts', []))
        total_tables += len(slide.get('tables', []))

    print(f"   📈 Charts: {type_counts.get('chart', 0)} slides ({total_charts} total charts)")
    print(f"   📊 Tables: {type_counts.get('table', 0)} slides ({total_tables} total tables)")
    print(f"   📝 Text: {type_counts.get('text', 0)}")
    print(f"   📌 Section headers: {type_counts.get('section_header', 0)}")

    return slides_data


def save_json(data, output_path):
    """Save data as formatted JSON."""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

    print(f"💾 Saved to: {output_path}")


def main():
    """Main entry point."""
    if len(sys.argv) < 2:
        print("Usage: python3 extract_ppt_v2.py <path_to_ppt>")
        print("Example: python3 extract_ppt_v2.py slide/survey-phase2-eng-PPT\\ \\(3\\).pptx")
        sys.exit(1)

    ppt_path = sys.argv[1]

    # Extract content
    slides_data = extract_presentation(ppt_path)

    # Save to JSON
    output_path = Path("output") / "extracted_slides_v2.json"
    save_json(slides_data, output_path)

    # Print sample
    print("\n📊 Sample Output:")
    for slide in slides_data[:2]:
        print(f"\n--- Slide {slide['id']} ({slide['type']}) ---")
        print(f"Title: {slide['title'][:60]}...")
        if slide['tables']:
            print(f"Tables: {len(slide['tables'])}")
            print(f"  First table: {len(slide['tables'][0]['headers'])} cols, {len(slide['tables'][0]['rows'])} data rows")
        else:
            preview = slide['content'][:80] + "..." if len(slide['content']) > 80 else slide['content']
            print(f"Content: {preview}")


if __name__ == "__main__":
    main()
