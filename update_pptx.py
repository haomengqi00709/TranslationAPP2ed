"""
Update PowerPoint file with aligned translated runs, tables, and charts
"""

import json
import logging
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from typing import Dict, Any, List
from pathlib import Path

logger = logging.getLogger(__name__)


class PowerPointUpdater:
    """Update PowerPoint presentations with translated content."""

    def __init__(self):
        """Initialize PowerPoint updater."""
        pass

    def update_presentation(
        self,
        input_pptx: str,
        aligned_jsonl: str,
        output_pptx: str
    ) -> Dict[str, int]:
        """
        Update PowerPoint presentation with aligned translations.

        Args:
            input_pptx: Path to input PowerPoint file
            aligned_jsonl: Path to JSONL file with aligned content (paragraphs, tables, charts)
            output_pptx: Path to output PowerPoint file

        Returns:
            Dict with counts: {"paragraphs": N, "tables": N, "charts": N}
        """
        logger.info(f"Updating presentation {input_pptx}")

        # Load presentation
        prs = Presentation(input_pptx)

        # Load aligned content
        content_items = self._load_content(aligned_jsonl)

        # Track counts by content type
        counts = {"paragraphs": 0, "tables": 0, "charts": 0}

        # Update each item based on content type
        for item in content_items:
            try:
                content_type = item.get("content_type", "paragraph")

                if content_type == "paragraph" or content_type == "text":
                    if self._update_paragraph(prs, item):
                        counts["paragraphs"] += 1

                elif content_type == "table":
                    if self._update_table(prs, item):
                        counts["tables"] += 1

                elif content_type == "chart":
                    if self._update_chart(prs, item):
                        counts["charts"] += 1

                else:
                    logger.warning(f"Unknown content type: {content_type}")

            except Exception as e:
                logger.error(f"Error updating {item.get('content_type', 'unknown')} "
                           f"on slide {item.get('slide_index', '?')}: {e}")
                import traceback
                traceback.print_exc()
                continue

        # Save presentation
        Path(output_pptx).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_pptx)
        logger.info(f"Updated {counts['paragraphs']} paragraphs, {counts['tables']} tables, "
                   f"{counts['charts']} charts, saved to {output_pptx}")

        return counts

    def _load_content(self, aligned_jsonl: str) -> List[Dict]:
        """Load aligned content (paragraphs, tables, charts) from JSONL file."""
        content = []
        with open(aligned_jsonl, 'r', encoding='utf-8') as f:
            for line in f:
                line = line.strip()
                if line:
                    try:
                        content.append(json.loads(line))
                    except json.JSONDecodeError as e:
                        logger.error(f"Error decoding JSON: {e}")
                        continue
        return content

    def _update_paragraph(self, prs: Presentation, para_data: Dict) -> bool:
        """
        Update a single paragraph with aligned runs.

        Returns:
            True if successful, False otherwise
        """
        try:
            slide_idx = para_data["slide_index"]
            shape_idx = para_data["shape_index"]
            para_idx = para_data["paragraph_index"]
            aligned_runs = para_data["aligned_runs"]

            # Get the paragraph
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]

            if not shape.has_text_frame:
                logger.warning(f"Shape {shape_idx} on slide {slide_idx} has no text frame")
                return False

            paragraph = shape.text_frame.paragraphs[para_idx]

            # Clear existing runs
            for run in paragraph.runs:
                run.text = ""

            # Add aligned runs
            for run_data in aligned_runs:
                run = paragraph.add_run()
                run.text = run_data["text"]

                # Apply formatting
                self._apply_run_formatting(run, run_data)

            return True

        except Exception as e:
            logger.error(f"Error updating paragraph: {e}")
            return False

    def _update_table(self, prs: Presentation, table_data: Dict) -> bool:
        """
        Update a table with translated cell content.

        Returns:
            True if successful, False otherwise
        """
        try:
            slide_idx = table_data["slide_index"]
            shape_idx = table_data["shape_index"]
            cells_data = table_data.get("cells", [])

            # Get the table
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]

            if not shape.has_table:
                logger.warning(f"Shape {shape_idx} on slide {slide_idx} is not a table")
                return False

            table = shape.table

            # Update each cell
            for cell_data in cells_data:
                row = cell_data["row"]
                col = cell_data["col"]
                cell = table.cell(row, col)

                # Update each paragraph in the cell
                for para_data in cell_data.get("paragraphs", []):
                    para_idx = para_data["paragraph_index"]

                    if para_idx >= len(cell.text_frame.paragraphs):
                        logger.warning(f"Paragraph index {para_idx} out of range for cell ({row}, {col})")
                        continue

                    paragraph = cell.text_frame.paragraphs[para_idx]

                    # Clear existing runs
                    for run in paragraph.runs:
                        run.text = ""

                    # Add aligned runs
                    for run_data in para_data.get("runs", []):
                        run = paragraph.add_run()
                        run.text = run_data["text"]

                        # Apply formatting
                        self._apply_run_formatting(run, run_data)

            logger.info(f"✓ Updated table on slide {slide_idx} with {len(cells_data)} cells")
            return True

        except Exception as e:
            logger.error(f"Error updating table: {e}")
            import traceback
            traceback.print_exc()
            return False

    def _update_chart(self, prs: Presentation, chart_data: Dict) -> bool:
        """
        Update a chart with translated titles, axes, legend entries, and category labels.

        Returns:
            True if successful, False otherwise
        """
        try:
            slide_idx = chart_data["slide_index"]
            shape_idx = chart_data["shape_index"]

            # Get the chart
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]

            if not isinstance(shape, GraphicFrame):
                logger.warning(f"Shape {shape_idx} on slide {slide_idx} is not a chart")
                return False

            chart = shape.chart

            # Update chart title (preserving formatting)
            if chart_data.get("title") and chart.has_title:
                title_info = chart_data["title"]
                translated_text = title_info["text"]

                # Clear existing text and add new with formatting
                text_frame = chart.chart_title.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = translated_text

                # Apply formatting from extracted data
                if title_info.get("font"):
                    run.font.name = title_info["font"]
                if title_info.get("size"):
                    run.font.size = Pt(float(title_info["size"]))
                if title_info.get("bold") is not None:
                    run.font.bold = bool(title_info["bold"])
                if title_info.get("italic") is not None:
                    run.font.italic = bool(title_info["italic"])
                if title_info.get("underline") is not None:
                    run.font.underline = bool(title_info["underline"])
                if title_info.get("color"):
                    self._apply_color(run, title_info["color"])

                logger.info(f"✓ Updated chart title: '{translated_text}' (with formatting)")

            # Update axis titles
            if chart_data.get("axis_titles", {}).get("category"):
                try:
                    if hasattr(chart, 'category_axis') and chart.category_axis:
                        if chart.category_axis.has_title:
                            translated_text = chart_data["axis_titles"]["category"]["text"]
                            chart.category_axis.axis_title.text_frame.text = translated_text
                            logger.info(f"✓ Updated X-axis title: '{translated_text}'")
                except Exception as e:
                    logger.warning(f"Error updating category axis: {e}")

            if chart_data.get("axis_titles", {}).get("value"):
                try:
                    if hasattr(chart, 'value_axis') and chart.value_axis:
                        if chart.value_axis.has_title:
                            translated_text = chart_data["axis_titles"]["value"]["text"]
                            chart.value_axis.axis_title.text_frame.text = translated_text
                            logger.info(f"✓ Updated Y-axis title: '{translated_text}'")
                except Exception as e:
                    logger.warning(f"Error updating value axis: {e}")

            # Update legend entries and category labels using chart.replace_data()
            legend_entries = chart_data.get("legend_entries", [])
            category_labels = chart_data.get("category_labels", [])

            if legend_entries or category_labels:
                try:
                    logger.info(f"Attempting to update chart data (legend + categories)")

                    # Get the plot to access series data
                    if hasattr(chart, 'plots') and chart.plots:
                        plot = chart.plots[0]

                        # Use data label settings from extracted data (if available)
                        data_label_settings = chart_data.get("data_label_settings", [])

                        # If not in extracted data, save current settings
                        if not data_label_settings and hasattr(plot, 'series'):
                            for series in plot.series:
                                if hasattr(series, 'data_labels') and series.data_labels:
                                    settings = {
                                        'series_index': len(data_label_settings),
                                        'show_value': series.data_labels.show_value if hasattr(series.data_labels, 'show_value') else False,
                                        'show_percentage': series.data_labels.show_percentage if hasattr(series.data_labels, 'show_percentage') else False,
                                        'show_category_name': series.data_labels.show_category_name if hasattr(series.data_labels, 'show_category_name') else False,
                                        'show_series_name': series.data_labels.show_series_name if hasattr(series.data_labels, 'show_series_name') else False,
                                        'number_format': series.data_labels.number_format if hasattr(series.data_labels, 'number_format') else None
                                    }
                                    data_label_settings.append(settings)
                                else:
                                    data_label_settings.append(None)

                        # Build new ChartData
                        new_chart_data = CategoryChartData()

                        # Add categories (translated)
                        if category_labels:
                            categories = [label["text"] for label in sorted(category_labels, key=lambda x: x["index"])]
                            new_chart_data.categories = categories
                            logger.info(f"  Categories: {categories}")
                        else:
                            # Keep original categories
                            if hasattr(plot, 'categories'):
                                new_chart_data.categories = list(plot.categories)

                        # Add series with translated names
                        if hasattr(plot, 'series'):
                            for series_idx, series in enumerate(plot.series):
                                # Find translated series name
                                series_name = series.name
                                for entry in legend_entries:
                                    if entry["series_index"] == series_idx:
                                        series_name = entry["text"]
                                        break

                                # Get values from original series
                                values = []
                                try:
                                    if hasattr(series, 'values'):
                                        values = list(series.values)
                                    elif hasattr(series, 'points'):
                                        values = [point.value for point in series.points if hasattr(point, 'value')]
                                except Exception as e:
                                    logger.warning(f"  Could not extract values from series {series_idx}: {e}")
                                    values = [0] * len(new_chart_data.categories)

                                new_chart_data.add_series(series_name, values)
                                logger.info(f"  Series {series_idx}: '{series_name}' with {len(values)} values")

                        # Replace chart data
                        chart.replace_data(new_chart_data)

                        # RESTORE data label settings after replace
                        if hasattr(chart, 'plots') and chart.plots:
                            plot = chart.plots[0]
                            if hasattr(plot, 'series'):
                                for series_idx, series in enumerate(plot.series):
                                    # Find settings for this series
                                    settings = None
                                    for s in data_label_settings:
                                        if s and s.get('series_index') == series_idx:
                                            settings = s
                                            break

                                    if settings:
                                        try:
                                            if hasattr(series, 'data_labels'):
                                                series.data_labels.show_value = settings['show_value']
                                                series.data_labels.show_percentage = settings['show_percentage']
                                                series.data_labels.show_category_name = settings['show_category_name']
                                                series.data_labels.show_series_name = settings['show_series_name']
                                                if settings.get('number_format'):
                                                    series.data_labels.number_format = settings['number_format']
                                                logger.info(f"  Restored data label format for series {series_idx}: percentage={settings['show_percentage']}, format={settings.get('number_format')}")
                                        except Exception as e:
                                            logger.warning(f"  Could not restore data labels for series {series_idx}: {e}")

                        logger.info(f"✓ Updated chart data (legend + categories)")

                except Exception as e:
                    logger.error(f"Error updating chart data: {e}")
                    import traceback
                    traceback.print_exc()

            return True

        except Exception as e:
            logger.error(f"Error updating chart: {e}")
            import traceback
            traceback.print_exc()
            return False

    def _apply_run_formatting(self, run, run_data: Dict):
        """Apply formatting to a run based on run_data."""
        if run.font:
            # Bold - explicitly set True or False
            if "bold" in run_data:
                run.font.bold = bool(run_data["bold"])

            # Italic - explicitly set True or False
            if "italic" in run_data:
                run.font.italic = bool(run_data["italic"])

            # Underline - explicitly set True or False
            if "underline" in run_data:
                run.font.underline = bool(run_data["underline"])

            # Font name
            if run_data.get("font"):
                run.font.name = run_data["font"]

            # Font size - MUST be set
            if run_data.get("size"):
                try:
                    run.font.size = Pt(float(run_data["size"]))
                except (ValueError, TypeError) as e:
                    logger.warning(f"Error setting font size {run_data.get('size')}: {e}")

            # Color - CRITICAL
            if run_data.get("color"):
                self._apply_color(run, run_data["color"])

            # Superscript
            if "superscript" in run_data:
                run.font.superscript = bool(run_data["superscript"])

            # Subscript
            if "subscript" in run_data:
                run.font.subscript = bool(run_data["subscript"])

        # Apply hyperlink
        if run_data.get("hyperlink"):
            try:
                run.hyperlink.address = run_data["hyperlink"]
            except Exception as e:
                logger.debug(f"Error setting hyperlink: {e}")

    def _apply_color(self, run, color_value: Any):
        """Apply color to run."""
        if not color_value:
            return

        try:
            # Handle hex color
            if isinstance(color_value, str) and color_value.startswith('#'):
                rgb_str = color_value.lstrip('#')
                if len(rgb_str) == 6:
                    r = int(rgb_str[0:2], 16)
                    g = int(rgb_str[2:4], 16)
                    b = int(rgb_str[4:6], 16)
                    run.font.color.rgb = RGBColor(r, g, b)

            # Handle theme color (theme:BACKGROUND_1 (14), etc.)
            elif isinstance(color_value, str) and color_value.startswith('theme:'):
                theme_str = color_value.replace('theme:', '').strip()
                # Extract the theme color name (e.g., "BACKGROUND_1" from "BACKGROUND_1 (14)")
                theme_name = theme_str.split('(')[0].strip()

                # Map theme names to MSO_THEME_COLOR enum
                theme_map = {
                    'ACCENT_1': MSO_THEME_COLOR.ACCENT_1,
                    'ACCENT_2': MSO_THEME_COLOR.ACCENT_2,
                    'ACCENT_3': MSO_THEME_COLOR.ACCENT_3,
                    'ACCENT_4': MSO_THEME_COLOR.ACCENT_4,
                    'ACCENT_5': MSO_THEME_COLOR.ACCENT_5,
                    'ACCENT_6': MSO_THEME_COLOR.ACCENT_6,
                    'BACKGROUND_1': MSO_THEME_COLOR.BACKGROUND_1,
                    'BACKGROUND_2': MSO_THEME_COLOR.BACKGROUND_2,
                    'DARK_1': MSO_THEME_COLOR.DARK_1,
                    'DARK_2': MSO_THEME_COLOR.DARK_2,
                    'LIGHT_1': MSO_THEME_COLOR.LIGHT_1,
                    'LIGHT_2': MSO_THEME_COLOR.LIGHT_2,
                    'TEXT_1': MSO_THEME_COLOR.TEXT_1,
                    'TEXT_2': MSO_THEME_COLOR.TEXT_2,
                    'HYPERLINK': MSO_THEME_COLOR.HYPERLINK,
                    'FOLLOWED_HYPERLINK': MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
                }

                if theme_name in theme_map:
                    run.font.color.theme_color = theme_map[theme_name]
                    logger.debug(f"Applied theme color: {theme_name}")
                else:
                    logger.warning(f"Unknown theme color: {theme_name}")

        except Exception as e:
            logger.warning(f"Error applying color {color_value}: {e}")


def main():
    """Example usage."""
    import sys

    if len(sys.argv) < 4:
        print("Usage: python update_pptx.py <input.pptx> <aligned.jsonl> <output.pptx>")
        sys.exit(1)

    updater = PowerPointUpdater()
    counts = updater.update_presentation(sys.argv[1], sys.argv[2], sys.argv[3])
    print(f"\n✅ Updated {counts['paragraphs']} paragraphs, {counts['tables']} tables, {counts['charts']} charts")


if __name__ == "__main__":
    main()
