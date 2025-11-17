"""
Comprehensive content extraction from PowerPoint files
Extracts text paragraphs, tables, and chart titles with full formatting
"""

import json
import logging
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import MSO_COLOR_TYPE
from typing import List, Dict, Any
from pathlib import Path

logger = logging.getLogger(__name__)


class ContentExtractor:
    """Extract text, tables, and charts from PowerPoint presentations."""

    def __init__(self):
        """Initialize content extractor."""
        pass

    def extract_all(
        self,
        pptx_path: str,
        text_output: str,
        table_output: str,
        chart_output: str
    ) -> Dict[str, int]:
        """
        Extract all content types from PowerPoint.

        Args:
            pptx_path: Path to input PowerPoint file
            text_output: Path to text paragraphs JSONL
            table_output: Path to tables JSONL
            chart_output: Path to chart titles JSONL

        Returns:
            Dictionary with counts of extracted items
        """
        logger.info(f"Extracting all content from {pptx_path}")

        text_count = self.extract_text_paragraphs(pptx_path, text_output)
        table_count = self.extract_tables(pptx_path, table_output)
        chart_count = self.extract_chart_titles(pptx_path, chart_output)

        return {
            "text_paragraphs": text_count,
            "tables": table_count,
            "chart_titles": chart_count
        }

    def extract_text_paragraphs(self, pptx_path: str, output_jsonl: str) -> int:
        """
        Extract text paragraphs (non-table, non-chart text).

        Args:
            pptx_path: Path to input PowerPoint file
            output_jsonl: Path to output JSONL file

        Returns:
            Number of paragraphs extracted
        """
        logger.info(f"Extracting text paragraphs from {pptx_path}")
        prs = Presentation(pptx_path)
        paragraphs = []

        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                # Skip tables and charts
                if shape.has_table or isinstance(shape, GraphicFrame):
                    continue

                if not shape.has_text_frame:
                    continue

                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                    # Skip empty paragraphs
                    if not paragraph.text.strip():
                        continue

                    # Extract paragraph-level info
                    paragraph_data = {
                        "content_type": "text",
                        "slide_index": slide_index,
                        "shape_index": shape_index,
                        "paragraph_index": paragraph_index,
                        "text": paragraph.text,  # Full paragraph text
                        "alignment": self._get_alignment(paragraph.alignment),
                        "level": paragraph.level,
                        "is_bullet": self._is_bullet(paragraph),
                        "runs": []
                    }

                    # Extract all runs within this paragraph
                    for run_index, run in enumerate(paragraph.runs):
                        run_data = {
                            "run_index": run_index,
                            "text": run.text,
                            "font": run.font.name if run.font and run.font.name else None,
                            "size": run.font.size.pt if run.font and run.font.size else None,
                            "bold": bool(run.font.bold) if run.font and run.font.bold is not None else False,
                            "italic": bool(run.font.italic) if run.font and run.font.italic is not None else False,
                            "underline": bool(run.font.underline) if run.font and run.font.underline is not None else False,
                            "color": self._get_color(run, paragraph),
                            "superscript": bool(run.font.superscript) if run.font and hasattr(run.font, 'superscript') and run.font.superscript is not None else False,
                            "subscript": bool(run.font.subscript) if run.font and hasattr(run.font, 'subscript') and run.font.subscript is not None else False,
                            "hyperlink": run.hyperlink.address if run.hyperlink and run.hyperlink.address else None
                        }

                        paragraph_data["runs"].append(run_data)

                    paragraphs.append(paragraph_data)

        # Write to JSONL file
        Path(output_jsonl).parent.mkdir(parents=True, exist_ok=True)
        with open(output_jsonl, 'w', encoding='utf-8') as f:
            for para in paragraphs:
                f.write(json.dumps(para, ensure_ascii=False) + '\n')

        logger.info(f"Extracted {len(paragraphs)} text paragraphs to {output_jsonl}")
        return len(paragraphs)

    def extract_tables(self, pptx_path: str, output_jsonl: str) -> int:
        """
        Extract tables with full cell/paragraph/run structure.

        Args:
            pptx_path: Path to input PowerPoint file
            output_jsonl: Path to output JSONL file

        Returns:
            Number of tables extracted

        Output format (one table per line):
        {
            "content_type": "table",
            "slide_index": 0,
            "shape_index": 1,
            "rows": 3,
            "cols": 2,
            "cells": [
                {
                    "row": 0,
                    "col": 0,
                    "paragraphs": [
                        {
                            "paragraph_index": 0,
                            "text": "Header 1",
                            "alignment": "left",
                            "level": 0,
                            "is_bullet": false,
                            "runs": [...]
                        }
                    ]
                },
                ...
            ]
        }
        """
        logger.info(f"Extracting tables from {pptx_path}")
        prs = Presentation(pptx_path)
        tables = []

        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if not shape.has_table:
                    continue

                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)

                table_data = {
                    "content_type": "table",
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "rows": rows,
                    "cols": cols,
                    "cells": []
                }

                # Extract each cell
                for r in range(rows):
                    for c in range(cols):
                        cell = table.cell(r, c)
                        cell_data = {
                            "row": r,
                            "col": c,
                            "paragraphs": []
                        }

                        if cell.text_frame:
                            for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                                # Skip empty paragraphs
                                if not paragraph.text.strip():
                                    continue

                                para_data = {
                                    "paragraph_index": para_idx,
                                    "text": paragraph.text,
                                    "alignment": self._get_alignment(paragraph.alignment),
                                    "level": paragraph.level,
                                    "is_bullet": self._is_bullet(paragraph),
                                    "runs": []
                                }

                                # Extract runs with ordinal detection
                                previous_run_text = None
                                for run_idx, run in enumerate(paragraph.runs):
                                    run_data = {
                                        "run_index": run_idx,
                                        "text": run.text,
                                        "font": run.font.name if run.font and run.font.name else None,
                                        "size": run.font.size.pt if run.font and run.font.size else None,
                                        "bold": bool(run.font.bold) if run.font and run.font.bold is not None else False,
                                        "italic": bool(run.font.italic) if run.font and run.font.italic is not None else False,
                                        "underline": bool(run.font.underline) if run.font and run.font.underline is not None else False,
                                        "color": self._get_color(run, paragraph),
                                        "hyperlink": run.hyperlink.address if run.hyperlink and run.hyperlink.address else None
                                    }

                                    # Detect superscript (ordinal suffix heuristic)
                                    is_ordinal_suffix = run.text.strip().lower() in ['th', 'st', 'nd', 'rd']
                                    follows_number = previous_run_text is not None and previous_run_text.replace(',', '').strip().isdigit()
                                    superscript_value = bool(run.font.superscript) if run.font and hasattr(run.font, 'superscript') and run.font.superscript is not None else False
                                    run_data["superscript"] = superscript_value or (is_ordinal_suffix and follows_number)

                                    # Detect subscript (chemical formula heuristic)
                                    is_subscript_candidate = (
                                        previous_run_text is not None and
                                        previous_run_text.strip().isalpha() and
                                        run.text.strip().isdigit() and
                                        len(run.text.strip()) == 1
                                    )
                                    subscript_value = bool(run.font.subscript) if run.font and hasattr(run.font, 'subscript') and run.font.subscript is not None else False
                                    run_data["subscript"] = subscript_value or is_subscript_candidate

                                    para_data["runs"].append(run_data)
                                    previous_run_text = run.text

                                cell_data["paragraphs"].append(para_data)

                        table_data["cells"].append(cell_data)

                tables.append(table_data)

        # Write to JSONL file
        Path(output_jsonl).parent.mkdir(parents=True, exist_ok=True)
        with open(output_jsonl, 'w', encoding='utf-8') as f:
            for table in tables:
                f.write(json.dumps(table, ensure_ascii=False) + '\n')

        logger.info(f"Extracted {len(tables)} tables to {output_jsonl}")
        return len(tables)

    def extract_chart_titles(self, pptx_path: str, output_jsonl: str) -> int:
        """
        Extract comprehensive chart data: titles, axis labels, legend, categories, data labels.

        Args:
            pptx_path: Path to input PowerPoint file
            output_jsonl: Path to output JSONL file

        Returns:
            Number of charts extracted

        Output format (one chart per line):
        {
            "content_type": "chart",
            "slide_index": 0,
            "shape_index": 2,
            "chart_type": "PIE",
            "title": {
                "text": "Sales by Region",
                "font": "Calibri",
                "size": 18.0,
                "bold": true,
                "color": "theme:ACCENT_1 (5)"
            },
            "axis_titles": {
                "category": {"text": "Quarter", ...},
                "value": {"text": "Revenue ($M)", ...}
            },
            "legend_entries": [
                {"text": "Sales", ...},
                {"text": "Profit", ...}
            ],
            "category_labels": [
                {"text": "Q1", ...},
                {"text": "Q2", ...}
            ],
            "data_labels": [
                {"series_index": 0, "point_index": 0, "text": "25%", ...}
            ]
        }
        """
        logger.info(f"Extracting charts from {pptx_path}")
        prs = Presentation(pptx_path)
        charts = []

        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if not isinstance(shape, GraphicFrame):
                    continue

                try:
                    chart = shape.chart
                except (ValueError, AttributeError):
                    continue

                chart_type = str(getattr(chart, 'chart_type', 'Unknown'))

                chart_data = {
                    "content_type": "chart",
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "chart_type": chart_type,
                    "title": None,
                    "axis_titles": {},
                    "legend_entries": [],
                    "category_labels": [],
                    "data_labels": [],
                    "data_label_settings": []  # Format settings per series
                }

                # Extract chart title
                if chart.has_title:
                    title_text = chart.chart_title.text_frame.text.strip()
                    if title_text:
                        chart_data["title"] = self._extract_text_frame_formatting(
                            chart.chart_title.text_frame
                        )

                # Extract axis titles
                try:
                    # Category axis (X-axis)
                    if hasattr(chart, 'category_axis') and chart.category_axis:
                        if chart.category_axis.has_title:
                            axis_title_text = chart.category_axis.axis_title.text_frame.text.strip()
                            if axis_title_text:
                                chart_data["axis_titles"]["category"] = self._extract_text_frame_formatting(
                                    chart.category_axis.axis_title.text_frame
                                )

                    # Value axis (Y-axis)
                    if hasattr(chart, 'value_axis') and chart.value_axis:
                        if chart.value_axis.has_title:
                            axis_title_text = chart.value_axis.axis_title.text_frame.text.strip()
                            if axis_title_text:
                                chart_data["axis_titles"]["value"] = self._extract_text_frame_formatting(
                                    chart.value_axis.axis_title.text_frame
                                )
                except Exception as e:
                    logger.debug(f"Error extracting axis titles: {e}")

                # Extract legend entries
                try:
                    if hasattr(chart, 'plots') and chart.plots:
                        for plot_idx, plot in enumerate(chart.plots):
                            if hasattr(plot, 'series'):
                                for series_idx, series in enumerate(plot.series):
                                    if series.name:
                                        chart_data["legend_entries"].append({
                                            "series_index": series_idx,
                                            "text": str(series.name)
                                        })
                except Exception as e:
                    logger.debug(f"Error extracting legend entries: {e}")

                # Extract category labels (X-axis labels like "Q1", "Q2", etc.)
                try:
                    if hasattr(chart, 'plots') and chart.plots:
                        plot = chart.plots[0]
                        if hasattr(plot, 'categories') and plot.categories:
                            for cat_idx, category in enumerate(plot.categories):
                                if category:
                                    chart_data["category_labels"].append({
                                        "index": cat_idx,
                                        "text": str(category)
                                    })
                except Exception as e:
                    logger.debug(f"Error extracting category labels: {e}")

                # Extract data labels (values shown on chart points/bars) and their format settings
                try:
                    if hasattr(chart, 'plots') and chart.plots:
                        for plot_idx, plot in enumerate(chart.plots):
                            if hasattr(plot, 'series'):
                                for series_idx, series in enumerate(plot.series):
                                    if hasattr(series, 'data_labels') and series.data_labels:
                                        try:
                                            # Extract data label format settings
                                            settings = {
                                                'series_index': series_idx,
                                                'show_value': series.data_labels.show_value if hasattr(series.data_labels, 'show_value') else False,
                                                'show_percentage': series.data_labels.show_percentage if hasattr(series.data_labels, 'show_percentage') else False,
                                                'show_category_name': series.data_labels.show_category_name if hasattr(series.data_labels, 'show_category_name') else False,
                                                'show_series_name': series.data_labels.show_series_name if hasattr(series.data_labels, 'show_series_name') else False,
                                                'number_format': series.data_labels.number_format if hasattr(series.data_labels, 'number_format') else None
                                            }
                                            chart_data["data_label_settings"].append(settings)

                                            # Extract data label text (if visible)
                                            if series.data_labels.show_value or series.data_labels.show_percentage:
                                                if hasattr(series, 'points'):
                                                    for point_idx, point in enumerate(series.points):
                                                        if hasattr(point, 'data_label') and point.data_label:
                                                            try:
                                                                label_text = point.data_label.text_frame.text.strip()
                                                                if label_text:
                                                                    chart_data["data_labels"].append({
                                                                        "series_index": series_idx,
                                                                        "point_index": point_idx,
                                                                        "text": label_text
                                                                    })
                                                            except:
                                                                pass
                                        except Exception as e:
                                            logger.debug(f"Error extracting data labels for series {series_idx}: {e}")
                except Exception as e:
                    logger.debug(f"Error extracting data labels: {e}")

                charts.append(chart_data)

        # Write to JSONL file
        Path(output_jsonl).parent.mkdir(parents=True, exist_ok=True)
        with open(output_jsonl, 'w', encoding='utf-8') as f:
            for chart in charts:
                f.write(json.dumps(chart, ensure_ascii=False) + '\n')

        logger.info(f"Extracted {len(charts)} charts to {output_jsonl}")
        return len(charts)

    def _extract_text_frame_formatting(self, text_frame) -> Dict[str, Any]:
        """
        Extract text and formatting from a text frame.
        Used for chart titles and axis labels.
        """
        full_text = text_frame.text.strip()
        font = size = bold = italic = underline = color = None

        # Get formatting from first run with values
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if not font and run.font and run.font.name:
                    font = run.font.name
                if not size and run.font and run.font.size:
                    size = run.font.size.pt
                if bold is None and run.font and run.font.bold is not None:
                    bold = run.font.bold
                if italic is None and run.font and run.font.italic is not None:
                    italic = run.font.italic
                if underline is None and run.font and run.font.underline is not None:
                    underline = run.font.underline
                if not color:
                    color = self._get_color(run, paragraph)

            # Fallback to paragraph font
            if not font and paragraph.font and paragraph.font.name:
                font = paragraph.font.name
            if not size and paragraph.font and paragraph.font.size:
                size = paragraph.font.size.pt
            if bold is None and paragraph.font and paragraph.font.bold is not None:
                bold = paragraph.font.bold
            if italic is None and paragraph.font and paragraph.font.italic is not None:
                italic = paragraph.font.italic
            if underline is None and paragraph.font and paragraph.font.underline is not None:
                underline = paragraph.font.underline

        return {
            "text": full_text,
            "font": font,
            "size": size,
            "bold": bold,
            "italic": italic,
            "underline": underline,
            "color": color
        }

    def _get_alignment(self, alignment) -> str:
        """Get alignment as string."""
        if alignment is None:
            return "left"

        alignment_map = {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify",
            PP_ALIGN.DISTRIBUTE: "distribute",
        }
        return alignment_map.get(alignment, "left")

    def _is_bullet(self, paragraph) -> bool:
        """Check if paragraph is a bullet point."""
        try:
            p = paragraph._p
            if p.pPr is not None:
                xml = p.pPr.xml
                if 'buChar' in xml or 'buAutoNum' in xml or 'buBlip' in xml:
                    return True
        except Exception as e:
            logger.debug(f"Error checking bullet status: {e}")
        return False

    def _get_color(self, run, paragraph) -> Any:
        """
        Get color information from run, with paragraph fallback.
        Returns format: "theme:ACCENT_1 (5)" or "#FF0000"
        """
        def extract_color(color_obj):
            if not color_obj:
                return None
            try:
                # Check for RGB color
                if color_obj.type == MSO_COLOR_TYPE.RGB and color_obj.rgb:
                    return f"#{color_obj.rgb}"

                # Check for theme color
                if color_obj.type == MSO_COLOR_TYPE.SCHEME and color_obj.theme_color:
                    theme_name = color_obj.theme_color.name
                    theme_idx = color_obj.theme_color.value
                    return f"theme:{theme_name} ({theme_idx})"

                return None
            except Exception as e:
                logger.debug(f"Error extracting color: {e}")
                return None

        # Try run-level color first
        if run.font and run.font.color:
            result = extract_color(run.font.color)
            if result:
                return result

        # Fallback to paragraph-level color
        if paragraph.font and paragraph.font.color:
            result = extract_color(paragraph.font.color)
            if result:
                return result

        return None


def main():
    """Example usage."""
    import sys

    if len(sys.argv) < 2:
        print("Usage: python extract_content.py <input.pptx>")
        sys.exit(1)

    extractor = ContentExtractor()
    counts = extractor.extract_all(
        pptx_path=sys.argv[1],
        text_output="temp/extracted_text.jsonl",
        table_output="temp/extracted_tables.jsonl",
        chart_output="temp/extracted_charts.jsonl"
    )

    print(f"Extracted:")
    print(f"  - {counts['text_paragraphs']} text paragraphs")
    print(f"  - {counts['tables']} tables")
    print(f"  - {counts['chart_titles']} chart titles")


if __name__ == "__main__":
    main()
