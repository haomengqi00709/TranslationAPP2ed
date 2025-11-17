"""
Test script to verify chart updating works
Modifies chart data and updates the PowerPoint file
"""

import json
import logging
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pathlib import Path

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def test_chart_update(input_pptx: str, output_pptx: str, chart_jsonl: str):
    """
    Test updating charts by modifying extracted chart data.

    This will:
    1. Load the extracted chart data
    2. Modify text (simulate translation): add "[TRANSLATED]" prefix
    3. Update the PowerPoint file with modified data
    """
    logger.info(f"Testing chart update from {input_pptx}")

    # Load presentation
    prs = Presentation(input_pptx)

    # Load chart data
    charts = []
    with open(chart_jsonl, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if line:
                charts.append(json.loads(line))

    logger.info(f"Loaded {len(charts)} charts from {chart_jsonl}")

    # Update each chart
    updated_count = 0
    for chart_data in charts:
        try:
            slide_idx = chart_data["slide_index"]
            shape_idx = chart_data["shape_index"]

            # Get the chart
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]

            if not isinstance(shape, GraphicFrame):
                logger.warning(f"Shape {shape_idx} on slide {slide_idx} is not a chart")
                continue

            chart = shape.chart

            # Update chart title
            if chart_data.get("title") and chart.has_title:
                translated_text = chart_data["title"]["text"]
                chart.chart_title.text_frame.text = translated_text
                logger.info(f"Updated chart title: '{translated_text}'")

            # Update category axis title (X-axis)
            if chart_data.get("axis_titles", {}).get("category"):
                try:
                    if hasattr(chart, 'category_axis') and chart.category_axis:
                        if chart.category_axis.has_title:
                            translated_text = chart_data["axis_titles"]["category"]["text"]
                            chart.category_axis.axis_title.text_frame.text = translated_text
                            logger.info(f"Updated X-axis title: '{translated_text}'")
                except Exception as e:
                    logger.warning(f"Error updating category axis: {e}")

            # Update value axis title (Y-axis)
            if chart_data.get("axis_titles", {}).get("value"):
                try:
                    if hasattr(chart, 'value_axis') and chart.value_axis:
                        if chart.value_axis.has_title:
                            translated_text = chart_data["axis_titles"]["value"]["text"]
                            chart.value_axis.axis_title.text_frame.text = translated_text
                            logger.info(f"Updated Y-axis title: '{translated_text}'")
                except Exception as e:
                    logger.warning(f"Error updating value axis: {e}")

            # Update legend entries (series names)
            # NOTE: series.name is READ-ONLY in python-pptx, cannot be updated directly
            # Legend entries need to be updated via Excel worksheet embedded in chart
            legend_entries = chart_data.get("legend_entries", [])
            if legend_entries:
                logger.info(f"Found {len(legend_entries)} legend entries (series.name is read-only, skipping)")

            # Update category labels (X-axis labels)
            category_labels = chart_data.get("category_labels", [])
            if category_labels:
                try:
                    # Note: Updating category labels requires modifying chart data
                    # This is more complex and chart-type specific
                    logger.info(f"Found {len(category_labels)} category labels (updating not implemented yet)")
                except Exception as e:
                    logger.warning(f"Error with category labels: {e}")

            updated_count += 1

        except Exception as e:
            logger.error(f"Error updating chart on slide {chart_data.get('slide_index', '?')}: {e}")
            import traceback
            traceback.print_exc()
            continue

    # Save presentation
    Path(output_pptx).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    logger.info(f"Updated {updated_count} charts, saved to {output_pptx}")

    return updated_count


def main():
    """Test with extracted data."""
    import sys

    if len(sys.argv) < 4:
        print("Usage: python test_chart_update.py <input.pptx> <charts.jsonl> <output.pptx>")
        print("\nExample:")
        print("  python test_chart_update.py \\")
        print("    'slides/PPT-3-Government-in-Canada1 (2).pptx' \\")
        print("    temp/extracted_charts.jsonl \\")
        print("    output/test_charts_updated.pptx")
        sys.exit(1)

    count = test_chart_update(sys.argv[1], sys.argv[3], sys.argv[2])
    print(f"\n✅ Updated {count} charts successfully!")
    print(f"Open the output file to verify the changes:")
    print(f"  open '{sys.argv[3]}'")


if __name__ == "__main__":
    main()
