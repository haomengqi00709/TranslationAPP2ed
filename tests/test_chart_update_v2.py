"""
Test script v2 - Update charts using chart.replace_data()
This properly updates legend entries and category labels
"""

import json
import logging
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.chart.data import ChartData, CategoryChartData
from pptx.util import Pt
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pathlib import Path

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def apply_color(run, color_value):
    """Apply color to run (hex or theme color)."""
    if not color_value:
        return

    try:
        # Handle hex color
        if isinstance(color_value, str) and color_value.startswith('#'):
            rgb_str = color_value.lstrip('#')
            if len(rgb_str) == 6:
                r = int(rgb_str[0:2], 16)
                g = int(rgb_str[2:4], 16)
                b = int(rgb_str[4:6], 16)
                run.font.color.rgb = RGBColor(r, g, b)

        # Handle theme color (theme:ACCENT_1 (5), etc.)
        elif isinstance(color_value, str) and color_value.startswith('theme:'):
            theme_str = color_value.replace('theme:', '').strip()
            theme_name = theme_str.split('(')[0].strip()

            theme_map = {
                'ACCENT_1': MSO_THEME_COLOR.ACCENT_1,
                'ACCENT_2': MSO_THEME_COLOR.ACCENT_2,
                'ACCENT_3': MSO_THEME_COLOR.ACCENT_3,
                'ACCENT_4': MSO_THEME_COLOR.ACCENT_4,
                'ACCENT_5': MSO_THEME_COLOR.ACCENT_5,
                'ACCENT_6': MSO_THEME_COLOR.ACCENT_6,
                'BACKGROUND_1': MSO_THEME_COLOR.BACKGROUND_1,
                'BACKGROUND_2': MSO_THEME_COLOR.BACKGROUND_2,
                'DARK_1': MSO_THEME_COLOR.DARK_1,
                'DARK_2': MSO_THEME_COLOR.DARK_2,
                'LIGHT_1': MSO_THEME_COLOR.LIGHT_1,
                'LIGHT_2': MSO_THEME_COLOR.LIGHT_2,
                'TEXT_1': MSO_THEME_COLOR.TEXT_1,
                'TEXT_2': MSO_THEME_COLOR.TEXT_2,
            }

            if theme_name in theme_map:
                run.font.color.theme_color = theme_map[theme_name]
    except Exception as e:
        logger.warning(f"Error applying color {color_value}: {e}")


def test_chart_update_v2(input_pptx: str, output_pptx: str, chart_jsonl: str):
    """
    Test updating charts using chart.replace_data() method.
    This properly updates categories and series names (legend entries).
    """
    logger.info(f"Testing chart update v2 from {input_pptx}")

    # Load presentation
    prs = Presentation(input_pptx)

    # Load chart data
    charts = []
    with open(chart_jsonl, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if line:
                charts.append(json.loads(line))

    logger.info(f"Loaded {len(charts)} charts from {chart_jsonl}")

    # Update each chart
    updated_count = 0
    for chart_data in charts:
        try:
            slide_idx = chart_data["slide_index"]
            shape_idx = chart_data["shape_index"]

            # Get the chart
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]

            if not isinstance(shape, GraphicFrame):
                logger.warning(f"Shape {shape_idx} on slide {slide_idx} is not a chart")
                continue

            chart = shape.chart

            # Update chart title (preserving formatting)
            if chart_data.get("title") and chart.has_title:
                title_info = chart_data["title"]
                translated_text = title_info["text"]

                # Clear existing text and add new with formatting
                text_frame = chart.chart_title.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = translated_text

                # Apply formatting from extracted data
                if title_info.get("font"):
                    run.font.name = title_info["font"]
                if title_info.get("size"):
                    from pptx.util import Pt
                    run.font.size = Pt(float(title_info["size"]))
                if title_info.get("bold") is not None:
                    run.font.bold = bool(title_info["bold"])
                if title_info.get("italic") is not None:
                    run.font.italic = bool(title_info["italic"])
                if title_info.get("underline") is not None:
                    run.font.underline = bool(title_info["underline"])
                if title_info.get("color"):
                    apply_color(run, title_info["color"])

                logger.info(f"✓ Updated chart title: '{translated_text}' (with formatting)")

            # Update axis titles
            if chart_data.get("axis_titles", {}).get("category"):
                try:
                    if hasattr(chart, 'category_axis') and chart.category_axis:
                        if chart.category_axis.has_title:
                            translated_text = chart_data["axis_titles"]["category"]["text"]
                            chart.category_axis.axis_title.text_frame.text = translated_text
                            logger.info(f"✓ Updated X-axis title: '{translated_text}'")
                except Exception as e:
                    logger.warning(f"Error updating category axis: {e}")

            if chart_data.get("axis_titles", {}).get("value"):
                try:
                    if hasattr(chart, 'value_axis') and chart.value_axis:
                        if chart.value_axis.has_title:
                            translated_text = chart_data["axis_titles"]["value"]["text"]
                            chart.value_axis.axis_title.text_frame.text = translated_text
                            logger.info(f"✓ Updated Y-axis title: '{translated_text}'")
                except Exception as e:
                    logger.warning(f"Error updating value axis: {e}")

            # Update legend entries and category labels using chart.replace_data()
            legend_entries = chart_data.get("legend_entries", [])
            category_labels = chart_data.get("category_labels", [])

            if legend_entries or category_labels:
                try:
                    # Extract current chart data
                    logger.info(f"Attempting to update chart data (legend + categories)")

                    # Get the plot to access series data
                    if hasattr(chart, 'plots') and chart.plots:
                        plot = chart.plots[0]

                        # Use data label settings from extracted data (if available)
                        data_label_settings = chart_data.get("data_label_settings", [])

                        # If not in extracted data, save current settings
                        if not data_label_settings and hasattr(plot, 'series'):
                            for series in plot.series:
                                if hasattr(series, 'data_labels') and series.data_labels:
                                    settings = {
                                        'series_index': len(data_label_settings),
                                        'show_value': series.data_labels.show_value if hasattr(series.data_labels, 'show_value') else False,
                                        'show_percentage': series.data_labels.show_percentage if hasattr(series.data_labels, 'show_percentage') else False,
                                        'show_category_name': series.data_labels.show_category_name if hasattr(series.data_labels, 'show_category_name') else False,
                                        'show_series_name': series.data_labels.show_series_name if hasattr(series.data_labels, 'show_series_name') else False,
                                        'number_format': series.data_labels.number_format if hasattr(series.data_labels, 'number_format') else None
                                    }
                                    data_label_settings.append(settings)
                                else:
                                    data_label_settings.append(None)

                        # Build new ChartData
                        new_chart_data = CategoryChartData()

                        # Add categories (translated)
                        if category_labels:
                            categories = [label["text"] for label in sorted(category_labels, key=lambda x: x["index"])]
                            new_chart_data.categories = categories
                            logger.info(f"  Categories: {categories}")
                        else:
                            # Keep original categories
                            if hasattr(plot, 'categories'):
                                new_chart_data.categories = list(plot.categories)

                        # Add series with translated names
                        if hasattr(plot, 'series'):
                            for series_idx, series in enumerate(plot.series):
                                # Find translated series name
                                series_name = series.name
                                for entry in legend_entries:
                                    if entry["series_index"] == series_idx:
                                        series_name = entry["text"]
                                        break

                                # Get values from original series
                                values = []
                                try:
                                    if hasattr(series, 'values'):
                                        values = list(series.values)
                                    elif hasattr(series, 'points'):
                                        values = [point.value for point in series.points if hasattr(point, 'value')]
                                except Exception as e:
                                    logger.warning(f"  Could not extract values from series {series_idx}: {e}")
                                    values = [0] * len(new_chart_data.categories)

                                new_chart_data.add_series(series_name, values)
                                logger.info(f"  Series {series_idx}: '{series_name}' with {len(values)} values")

                        # Replace chart data
                        chart.replace_data(new_chart_data)

                        # RESTORE data label settings after replace
                        if hasattr(chart, 'plots') and chart.plots:
                            plot = chart.plots[0]
                            if hasattr(plot, 'series'):
                                for series_idx, series in enumerate(plot.series):
                                    # Find settings for this series
                                    settings = None
                                    for s in data_label_settings:
                                        if s and s.get('series_index') == series_idx:
                                            settings = s
                                            break

                                    if settings:
                                        try:
                                            if hasattr(series, 'data_labels'):
                                                series.data_labels.show_value = settings['show_value']
                                                series.data_labels.show_percentage = settings['show_percentage']
                                                series.data_labels.show_category_name = settings['show_category_name']
                                                series.data_labels.show_series_name = settings['show_series_name']
                                                if settings.get('number_format'):
                                                    series.data_labels.number_format = settings['number_format']
                                                logger.info(f"  Restored data label format for series {series_idx}: percentage={settings['show_percentage']}, format={settings.get('number_format')}")
                                        except Exception as e:
                                            logger.warning(f"  Could not restore data labels for series {series_idx}: {e}")

                        logger.info(f"✓ Updated chart data (legend + categories)")

                except Exception as e:
                    logger.error(f"Error updating chart data: {e}")
                    import traceback
                    traceback.print_exc()

            updated_count += 1

        except Exception as e:
            logger.error(f"Error updating chart on slide {chart_data.get('slide_index', '?')}: {e}")
            import traceback
            traceback.print_exc()
            continue

    # Save presentation
    Path(output_pptx).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    logger.info(f"\n✅ Updated {updated_count} charts, saved to {output_pptx}")

    return updated_count


def main():
    """Test with extracted data."""
    import sys

    if len(sys.argv) < 4:
        print("Usage: python test_chart_update_v2.py <input.pptx> <charts.jsonl> <output.pptx>")
        print("\nExample:")
        print("  python test_chart_update_v2.py \\")
        print("    'slides/PPT-3-Government-in-Canada1 (2).pptx' \\")
        print("    temp/extracted_charts.jsonl \\")
        print("    output/test_charts_updated_v2.pptx")
        sys.exit(1)

    count = test_chart_update_v2(sys.argv[1], sys.argv[3], sys.argv[2])
    print(f"\n✅ Updated {count} charts successfully!")
    print(f"Open the output file to verify the changes:")
    print(f"  open '{sys.argv[3]}'")


if __name__ == "__main__":
    main()
